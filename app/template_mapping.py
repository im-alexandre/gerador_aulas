from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable

from pptx import Presentation


DEFAULT_WANTED = {
    "title": ["title", "subtitle"],
    "standard": ["title", "pip", "bullets", "image"],
    "code": ["title", "pip", "code", "bullets"],
}


def build_mapping_from_existing_slides(
    pptx_path: Path,
    wanted: dict[str, list[str]] | None = None,
) -> dict:
    """Gera mapping de placeholder idx a partir de slides existentes."""
    wanted = wanted or DEFAULT_WANTED
    prs = Presentation(str(pptx_path))

    mapping = {"layouts": {}, "idx": {}}

    for slide in prs.slides:
        layout_name = slide.slide_layout.name.strip()
        if layout_name not in wanted:
            continue

        found: dict[str, int] = {}
        for ph in slide.placeholders:
            nm = (ph.name or "").strip()
            if nm in wanted[layout_name]:
                found[nm] = int(ph.placeholder_format.idx)

        if found:
            mapping["layouts"][layout_name] = layout_name
            mapping["idx"][layout_name] = found

    for layout_name, names in wanted.items():
        if layout_name not in mapping["idx"]:
            # fallback: buscar placeholders direto do layout
            layout = None
            for master in prs.slide_masters:
                for candidate in master.slide_layouts:
                    if candidate.name.strip() == layout_name:
                        layout = candidate
                        break
                if layout:
                    break
            if layout:
                found: dict[str, int] = {}
                for ph in layout.placeholders:
                    nm = (ph.name or "").strip()
                    if nm in names:
                        found[nm] = int(ph.placeholder_format.idx)
                if found:
                    mapping["layouts"][layout_name] = layout_name
                    mapping["idx"][layout_name] = found

        if layout_name not in mapping["idx"]:
            raise RuntimeError(
                f"Não achei slide exemplo para layout '{layout_name}'."
            )
        missing = [n for n in names if n not in mapping["idx"][layout_name]]
        if missing:
            raise RuntimeError(
                f"Layout '{layout_name}' sem placeholders: {missing}"
            )

    return mapping


def load_mapping(path: Path) -> dict:
    """Carrega mapping JSON de disco."""
    return json.loads(path.read_text(encoding="utf-8"))


def save_mapping(mapping: dict, path: Path) -> None:
    """Salva mapping JSON em disco."""
    path.write_text(json.dumps(mapping, ensure_ascii=False, indent=2), encoding="utf-8")


def map_path_for_template(pptx_path: Path) -> Path:
    """Retorna o path do mapping associado ao template."""
    return pptx_path.with_name(f"{pptx_path.stem}_map.json")


def ensure_template_mapping(
    pptx_path: Path,
    *,
    force: bool = False,
    wanted: dict[str, list[str]] | None = None,
) -> Path:
    """Gera o mapping se necessário e retorna o path."""
    map_path = map_path_for_template(pptx_path)
    if map_path.exists() and not force:
        return map_path

    mapping = build_mapping_from_existing_slides(pptx_path, wanted=wanted)
    save_mapping(mapping, map_path)
    return map_path


def validate_template_layouts(
    pptx_path: Path,
    *,
    wanted: dict[str, list[str]] | None = None,
) -> None:
    """Valida que o template contém os layouts e placeholders esperados."""
    build_mapping_from_existing_slides(pptx_path, wanted=wanted)


def load_or_build_mapping(
    pptx_path: Path,
    map_path: Path,
    wanted: dict[str, list[str]] | None = None,
) -> dict:
    """Deprecated: use ensure_template_mapping + load_mapping."""
    if map_path.exists():
        return load_mapping(map_path)
    mapping = build_mapping_from_existing_slides(pptx_path, wanted=wanted)
    save_mapping(mapping, map_path)
    return mapping
