from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _pick_main_picture(slide):
    pictures = [
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    if not pictures:
        return None
    return max(pictures, key=lambda sh: sh.width * sh.height)


def extract_slide_images(
    pptx_path: Path,
    out_dir: Path,
    slide_indices: list[int],
    slide_ids: list[str],
) -> list[Path | None]:
    """Extrai uma imagem por slide (maior) e salva em out_dir."""
    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(pptx_path))
    saved: list[Path | None] = []

    for idx, slide_id in zip(slide_indices, slide_ids):
        if idx < 0 or idx >= len(prs.slides):
            saved.append(None)
            continue
        slide = prs.slides[idx]
        pic = _pick_main_picture(slide)
        if not pic:
            saved.append(None)
            continue
        ext = (pic.image.ext or "png").lower()
        filename = f"gen_{slide_id}.{ext}"
        out_path = out_dir / filename
        out_path.write_bytes(pic.image.blob)
        saved.append(out_path)

    return saved
