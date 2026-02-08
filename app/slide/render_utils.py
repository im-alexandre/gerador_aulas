from __future__ import annotations

from pathlib import Path

from pptx.util import Pt


def get_placeholder_by_idx(slide, idx: int | None):
    """Busca placeholder pelo idx; retorna None se não existir."""
    if idx is None:
        return None
    try:
        return slide.placeholders[int(idx)]
    except (KeyError, ValueError, TypeError):
        return None


def set_text(shape, text: str) -> None:
    """Define texto em um shape com text_frame."""
    if not shape or not shape.has_text_frame:
        return
    shape.text_frame.text = text or ""


def set_bullets(shape, bullets: list[str]) -> None:
    """Preenche um shape de texto com bullets."""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.text = ""
    if not bullets:
        return
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0


def set_lead_with_bullets(shape, lead: str, bullets: list[str]) -> None:
    """Combina lead + bullets no mesmo shape."""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.text = ""
    if lead:
        tf.text = lead
    if bullets:
        if not lead:
            tf.text = bullets[0]
            start_idx = 1
        else:
            start_idx = 0
        for bullet in bullets[start_idx:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0


def _preserve_indent(line: str) -> str:
    if not line:
        return ""
    prefix = len(line) - len(line.lstrip(" "))
    if prefix <= 0:
        return line
    return ("\u00A0" * prefix) + line[prefix:]


def set_code(shape, code_text: str, font_name: str = "Consolas") -> None:
    """Preenche um shape com código monoespaçado, preservando identação."""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.text = ""
    lines = code_text.splitlines() or [""]
    for idx, line in enumerate(lines):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        run = p.add_run()
        text = _preserve_indent(line)
        run.text = text
        r = run._r
        t = r.find("{http://schemas.openxmlformats.org/drawingml/2006/main}t")
        if t is not None:
            t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        run.font.name = font_name
        run.font.size = Pt(16)


def resolve_image_path(base_dir: Path, image_path: str) -> Path:
    """Resolve o caminho da imagem a partir de uma base."""
    path = Path(image_path)
    if path.is_absolute():
        return path
    return base_dir / path
