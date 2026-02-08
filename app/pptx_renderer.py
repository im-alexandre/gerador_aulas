from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation

from app.slide import get_slide_class
from app.template_mapping import load_mapping, map_path_for_template


def find_layout_by_name(prs: Presentation, name: str):
    """Busca um layout pelo nome em todos os slide masters."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name.strip() == name.strip():
                return layout
    return None


def load_plan(plan_path: Path) -> dict[str, Any]:
    """Carrega o JSON do plano de slides."""
    return json.loads(plan_path.read_text(encoding="utf-8"))


def render_from_plan(
    plan: dict[str, Any],
    template_path: Path,
    output_path: Path,
    assets_base: Path,
    title: str | None = None,
) -> None:
    """Renderiza o PPTX final a partir do JSON e do template."""
    prs = Presentation(str(template_path))
    map_path = map_path_for_template(template_path)
    if not map_path.exists():
        raise SystemExit(
            f"Mapping não encontrado: {map_path.name}. Gere o mapping antes de renderizar."
        )
    mapping = load_mapping(map_path)

    default_layout = prs.slide_layouts[0]
    layout_title_name = mapping.get("layouts", {}).get("title", "title")
    layout_standard_name = mapping.get("layouts", {}).get("standard", "standard")
    layout_code_name = mapping.get("layouts", {}).get("code", "code")

    layout_title = find_layout_by_name(prs, layout_title_name) or default_layout
    layout_standard = find_layout_by_name(prs, layout_standard_name) or default_layout
    layout_code = find_layout_by_name(prs, layout_code_name) or layout_standard

    sentinel_layouts = set(mapping.get("layouts", {}).values())
    if sentinel_layouts:
        delete_sentinel_slides(prs, sentinel_layouts)

    if title:
        title_cls = get_slide_class("title")
        slide_capa = prs.slides.add_slide(layout_title)
        if title_cls:
            ph_map = mapping.get("idx", {}).get("title", {})
            title_cls.render({"title": title}, slide_capa, assets_base, ph_map)

    for slide in plan.get("slides", []):
        kind = slide.get("kind", "standard")
        slide_cls = get_slide_class(kind) or get_slide_class("standard")
        if not slide_cls:
            continue

        layout_name = slide_cls.LAYOUT_NAME
        if layout_name == "title":
            layout = layout_title
        elif layout_name == "code":
            layout = layout_code
        elif layout_name == "standard":
            layout = layout_standard
        elif layout_name:
            layout = find_layout_by_name(
                prs, mapping.get("layouts", {}).get(layout_name, layout_name)
            ) or layout_standard
        else:
            layout = layout_standard

        dst_slide = prs.slides.add_slide(layout)
        ph_map = mapping.get("idx", {}).get(layout_name, {})
        slide_cls.render(slide, dst_slide, assets_base, ph_map)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def delete_slide(prs: Presentation, slide_index: int) -> None:
    """Remove um slide pelo índice (inclui relacionamento)."""
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    sld_id = list(sld_id_lst)[slide_index]
    r_id = sld_id.rId
    sld_id_lst.remove(sld_id)
    prs.part.drop_rel(r_id)


def delete_sentinel_slides(prs: Presentation, sentinel_layouts: set[str]) -> None:
    """Remove todos os slides cujo layout é sentinela."""
    to_delete = [
        idx
        for idx, slide in enumerate(prs.slides)
        if slide.slide_layout.name.strip() in sentinel_layouts
    ]
    for idx in reversed(to_delete):
        delete_slide(prs, idx)
